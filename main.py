import csv
import json
import re

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE
from pptx.shapes.placeholder import Shape
from pptx.slide import Slide
from pptx.util import Inches


def fetch_page(url):
    """Fetches the HTML content of the given URL."""
    try:
        response = requests.get(url, headers={"User-Agent": "Mozilla/5.0"})
        response.raise_for_status()  # Raise an error for bad status codes
        response.encoding = "utf-8"
        return response.text
    except requests.exceptions.RequestException as e:
        print(f"Error fetching page: {e}")
        return None


def parse_html(html):
    """Parses the HTML content using BeautifulSoup."""
    return BeautifulSoup(html, "html.parser", from_encoding="latin-1")


def extract_page_title(soup: BeautifulSoup):
    """Extracts and cleans the page title."""
    title_tag = soup.find("title")
    if not title_tag:
        return ""

    title = title_tag.get_text().strip()
    return title


def find_jf_game(soup: BeautifulSoup):
    for div_container in soup.find_all("div", class_="level3"):
        script_tag = div_container.find("script", type="application/json")
        if not script_tag:
            continue
        try:
            whole_json = json.loads(script_tag.string)
        except json.JSONDecodeError:
            continue
        try:
            title: str = whole_json["x"]["layout"]["title"]["text"]
        except KeyError:
            continue
        if "José Figueiras" in title:
            return whole_json
    print("Did not find José Figueiras.")
    return whole_json


def extract_quiz_data(soup: BeautifulSoup):
    whole_json = find_jf_game(soup)
    if not whole_json:
        print("No matching div with target title found.")
        return []

    try:
        quiz_data = whole_json["x"]["data"]
    except KeyError as e:
        print(f"Error parsing JSON: {e}")
        return []

    parsed_data = []
    for part in quiz_data:
        part_data = []
        if "text" not in part:
            continue
        for row in part["text"]:
            soup_text = BeautifulSoup(row, "html.parser")
            parts = re.split(r"<br\s*/>", str(soup_text))
            if len(parts) < 2:
                continue
            player_tag = parts[0].strip()
            player_soup = BeautifulSoup(player_tag, "html.parser")
            b_tag = player_soup.find("b")
            if b_tag:
                full_name_and_team = b_tag.get_text(strip=True)
                player_name = full_name_and_team.split(" - ")[0]
            final_part = parts[-1].strip()
            final_soup = BeautifulSoup(final_part, "html.parser")
            points_tag = final_soup.find("b")

            if points_tag:
                points = points_tag.get_text(strip=True)
            theme_xt_xp = parts[1].strip()
            pattern = re.compile(r"^(.*?)\s*\(xT\s*=\s*([\d.]+),\s*xP\s*=\s*([\d.]+)\)")
            match = pattern.search(theme_xt_xp)
            if match:
                theme = match.group(1).strip()
                # Remove "Parte X " from theme if present
                theme = re.sub(r"^Parte\s\d\s", "", theme).strip()
                xt = float(match.group(2))
                xp = float(match.group(3))
            else:
                continue
            for br in soup_text.find_all("br"):
                br.replace_with(" ")
            # Extract question
            question_tag = soup_text.find("i")
            question = question_tag.get_text() if question_tag else ""
            # Extract answer
            answer_tag = soup_text.find("b", string="Resposta")
            answer = answer_tag.find_next("i").get_text() if answer_tag else ""
            part_data.append(
                {
                    "theme": theme,
                    "xT": xt,
                    "xP": xp,
                    "question": question,
                    "answer": answer,
                    "player": player_name,
                    "guessed": points == "2",
                }
            )
        parsed_data.append(part_data)
    return parsed_data


def sort_quiz_data(parsed_data):
    """Sorts each part of quiz data alphabetically by theme, except themes starting with 'Mystery Box' come last."""
    for part in parsed_data:
        part.sort(key=lambda x: (x["theme"].startswith("Mystery Box"), x["theme"]))
    return parsed_data


# Load CSV file and extract data
def load_csv(filename):
    with open(filename, newline="", encoding="utf-8") as csvfile:
        reader = csv.DictReader(csvfile)
        data = list(reader)
    return data


# Extract unique themes and sort them as specified
def get_sorted_themes(data):
    return sorted(set(row["theme"] for row in data), key=lambda x: (x.startswith("Mystery Box"), x))


# Create PowerPoint presentation
def create_ppt(data, output_file):
    left = top = width = height = Inches(0.75)
    prs = Presentation()

    # Generate index slide
    slide_layout = prs.slide_layouts[1]
    home_slide = prs.slides.add_slide(slide_layout)
    title = home_slide.shapes.title
    title.text = output_file

    themes = get_sorted_themes(data)
    theme_to_first_question_slide = {}

    # Generate slides for each question
    question_slides = []
    for row in data:
        slide = prs.slides.add_slide(slide_layout)
        # title
        title = slide.shapes.title
        title.text = row["theme"]
        # question
        body_shape: Shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        tf.text = row["question"]

        if row["theme"] not in theme_to_first_question_slide:
            theme_to_first_question_slide[row["theme"]] = slide
        question_slides.append(slide)

        slide = prs.slides.add_slide(slide_layout)
        # title
        title = slide.shapes.title
        title.text = row["theme"]
        # answer
        body_shape: Shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = row["answer"]
        # footer with xP
        footer = slide.shapes.add_textbox(left * 2, prs.slide_height - top, prs.slide_width - 4 * width, height)
        footer.text_frame.text = f"xP: {row['xP']}"
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string("c0c0c0")

    # Add hyperlinks from the index slide to the first question slide of each theme
    for i, theme in enumerate(themes, 1):
        theme_slide: Slide = theme_to_first_question_slide[theme]
        index_shape: Shape = home_slide.shapes.placeholders[1]
        rId = home_slide.part.relate_to(theme_slide.part, RELATIONSHIP_TYPE.SLIDE)
        p = index_shape.text_frame.add_paragraph()
        r = p.add_run()
        r.text = theme
        rPr = r._r.get_or_add_rPr()
        hlinkClick = rPr.add_hlinkClick(rId)
        hlinkClick.set("action", "ppaction://hlinksldjump")

    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        # all slides have button back to the index slide
        home_button = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ACTION_BUTTON_HOME, prs.slide_width - left, prs.slide_height - top, width, height
        )
        home_button.click_action.target_slide = home_slide

    # Save presentation
    prs.save(f"{output_file}.pptx")
    print(f"Presentation saved as {output_file}")


def print_json(content: dict, filename: str):
    with open(filename, "w") as json_file:
        json.dump(content, json_file, indent=2)


def main():
    urls = ["https://quizportugal.pt/sites/default/files/pictures/QNpt15_7.html#jam-37-x-38-jfg"]
    for url in urls:
        html_content = fetch_page(url)
        if html_content:
            soup = parse_html(html_content)
            page_title = extract_page_title(soup)
            quiz_data = extract_quiz_data(soup)
            sorted_data = sort_quiz_data(quiz_data)
            season, week = [int(n) for n in re.findall(r"\d+", page_title)]
            quiz = {"season": season, "week": week}
            parts = []
            for i, part_data in enumerate(sorted_data, 1):
                ppt_filename = f"{page_title} - Parte {i}"
                themes = [theme for theme in get_sorted_themes(part_data) if not theme.startswith("Mystery Box")]
                parts.append({"sequence": i, "themes": themes, "questions": part_data})
                create_ppt(part_data, ppt_filename)
            quiz["parts"] = parts
            print_json(quiz, f"{page_title}.json")


if __name__ == "__main__":
    main()
